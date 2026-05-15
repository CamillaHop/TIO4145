"""Pre-process raw PDFs and PPTX files in resources/<section_id>/ into
<basename>.parsed.json files alongside the sources. Idempotent: skips files
whose .parsed.json already exists and is newer than the source.

Usage:
    python scripts/parse_resources.py
    python scripts/parse_resources.py --section section_01
    python scripts/parse_resources.py --force
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from config_loader import REPO_ROOT, load_config, write_json_atomic

MAX_FILE_BYTES = 50 * 1024 * 1024  # 50MB — bigger files are usually scans


def parse_pdf(path: Path) -> dict:
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        sys.exit("ERROR: pdfplumber not installed. Run: pip install pdfplumber")

    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            pages.append({"page": i, "text": text.strip()})
    return {"source": path.name, "type": "pdf", "pages": pages}


def parse_pptx(path: Path) -> dict:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        sys.exit("ERROR: python-pptx not installed. Run: pip install python-pptx")

    pages = []
    pres = Presentation(str(path))
    for i, slide in enumerate(pres.slides, start=1):
        bits = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                bits.append(shape.text)
        pages.append({"page": i, "text": "\n".join(bits).strip()})
    return {"source": path.name, "type": "pptx", "pages": pages}


def needs_reparse(source: Path, force: bool) -> tuple[bool, Path]:
    out = source.parent / f"{source.stem}.parsed.json"
    if force or not out.exists():
        return True, out
    # If the existing .parsed.json was produced by Reducto, don't overwrite
    # it with a cheaper pdfplumber pass. Reducto output has top-level "chunks".
    try:
        existing = json.loads(out.read_text(encoding="utf-8"))
        if isinstance(existing, dict) and "chunks" in existing:
            return False, out
    except Exception:
        pass
    return out.stat().st_mtime < source.stat().st_mtime, out


def process_section(section_id: str, force: bool) -> tuple[int, int]:
    """Returns (parsed_count, skipped_count)."""
    folder = REPO_ROOT / "resources" / section_id
    if not folder.exists():
        print(f"  (no folder for {section_id}, skipping)")
        return 0, 0

    parsed = 0
    skipped = 0
    for path in sorted(folder.iterdir()):
        if not path.is_file():
            continue
        suffix = path.suffix.lower()
        if suffix not in (".pdf", ".pptx"):
            continue
        if path.stat().st_size > MAX_FILE_BYTES:
            print(
                f"  WARN: {path.name} is {path.stat().st_size // (1024*1024)}MB — too big, skipping. "
                "Pre-process externally (e.g. Reducto) and drop the result as a parsed.json.",
                file=sys.stderr,
            )
            continue

        do_parse, out = needs_reparse(path, force)
        if not do_parse:
            skipped += 1
            continue

        print(f"  parsing {path.name}…")
        try:
            if suffix == ".pdf":
                data = parse_pdf(path)
            else:
                data = parse_pptx(path)
        except Exception as e:
            print(f"  ERROR parsing {path.name}: {e}", file=sys.stderr)
            continue

        write_json_atomic(out, data)
        parsed += 1

    return parsed, skipped


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--section", help="parse only this section id (e.g. section_01)")
    ap.add_argument(
        "--force",
        action="store_true",
        help="re-parse files even if their .parsed.json is already up to date",
    )
    args = ap.parse_args()

    cfg = load_config()
    sections = [s["id"] for s in cfg.get("sections", [])]
    if args.section:
        if args.section not in sections:
            sys.exit(
                f"ERROR: section '{args.section}' is not in course_config.json. "
                f"Known sections: {', '.join(sections) or '(none)'}"
            )
        sections = [args.section]

    total_parsed = 0
    total_skipped = 0
    for sid in sections:
        print(f"▶ {sid}")
        p, s = process_section(sid, args.force)
        total_parsed += p
        total_skipped += s

    print(
        f"\nParsed {total_parsed} file(s) across {len(sections)} section(s) "
        f"(skipped {total_skipped} already up-to-date)."
    )


if __name__ == "__main__":
    main()
